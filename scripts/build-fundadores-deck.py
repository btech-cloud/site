#!/usr/bin/env python3
"""Gera PPTX dos Fundamentos BTECH para import no Canva (sem logotipos)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "assets" / "presentations" / "btech-fundamentos-fundadores.pptx"

# Paleta alinhada ao site (sem símbolos BTech)
INK = RGBColor(5, 5, 16)
INK_SOFT = RGBColor(244, 242, 255)
PANEL = RGBColor(26, 11, 64)
VIOLET = RGBColor(110, 30, 230)
BLUE = RGBColor(20, 60, 255)
CYAN = RGBColor(0, 225, 180)
MAGENTA = RGBColor(200, 0, 210)
ORANGE = RGBColor(255, 175, 60)
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(250, 251, 255)
MUTED = RGBColor(120, 115, 145)


def set_slide_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_accent_bar(slide, left=0, top=6.85, width=10, height=0.08, color=VIOLET) -> None:
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_gradient_bar(slide) -> None:
    w = 1.6
    colors = [VIOLET, BLUE, CYAN, ORANGE]
    x = 1.2
    for c in colors:
        shape = slide.shapes.add_shape(1, Inches(x), Inches(6.78), Inches(w), Inches(0.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = c
        shape.line.fill.background()
        x += w + 0.05


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size=28,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    italic=False,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return tf


def add_bullets(slide, left, top, width, height, items, *, size=20, color=INK, title=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(size + 4)
        p.font.bold = True
        p.font.color.rgb = color
        first = False
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)


def slide_title(prs, title, subtitle=None, *, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = INK if dark else LIGHT_BG
    set_slide_bg(slide, bg)
    tc = WHITE if dark else INK
    sc = INK_SOFT if dark else MUTED
    add_text_box(slide, 0.9, 2.2, 8.5, 1.2, title, size=40, bold=True, color=tc, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, 0.9, 3.3, 8.2, 1.5, subtitle, size=22, color=sc, align=PP_ALIGN.LEFT)
    if dark:
        add_gradient_bar(slide)
    else:
        add_accent_bar(slide, color=VIOLET)
    return slide


def slide_discussion(prs, question, hint):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PANEL)
    add_text_box(slide, 0.7, 0.5, 3, 0.5, "DISCUSSÃO", size=14, bold=True, color=CYAN)
    add_text_box(slide, 0.9, 2.0, 8.3, 2.5, question, size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text_box(slide, 0.9, 5.0, 8.0, 1.2, hint, size=18, color=INK_SOFT, italic=True)
    add_gradient_bar(slide)
    return slide


def slide_quote(prs, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, INK)
    text = "\n\n".join(lines)
    add_text_box(slide, 1.0, 2.3, 8.0, 3.0, text, size=32, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
    add_gradient_bar(slide)
    return slide


def slide_split_futures(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_text_box(slide, 0.9, 0.55, 8.5, 0.8, "Dois futuros possíveis", size=32, bold=True, color=INK)
    # left panel
    left = slide.shapes.add_shape(1, Inches(0.7), Inches(1.5), Inches(4.3), Inches(4.8))
    left.fill.solid()
    left.fill.fore_color.rgb = RGBColor(35, 25, 55)
    left.line.color.rgb = MAGENTA
    add_text_box(slide, 1.0, 1.7, 3.8, 0.6, "Concentração", size=22, bold=True, color=WHITE)
    add_bullets(
        slide,
        1.0,
        2.4,
        3.8,
        3.5,
        [
            "Riqueza e conhecimento concentram",
            "Decisões concentram",
            "Poucos controlam muito",
            "Autonomia diminui",
        ],
        size=17,
        color=INK_SOFT,
    )
    # right panel
    right = slide.shapes.add_shape(1, Inches(5.2), Inches(1.5), Inches(4.3), Inches(4.8))
    right.fill.solid()
    right.fill.fore_color.rgb = RGBColor(15, 40, 70)
    right.line.color.rgb = CYAN
    add_text_box(slide, 5.5, 1.7, 3.8, 0.6, "Ampliação", size=22, bold=True, color=WHITE)
    add_bullets(
        slide,
        5.5,
        2.4,
        3.8,
        3.5,
        [
            "Conhecimento distribuído",
            "Oportunidades ampliadas",
            "PMEs ganham capacidade",
            "Mais liberdade humana",
        ],
        size=17,
        color=INK_SOFT,
    )
    add_accent_bar(slide)
    return slide


def slide_examples(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_text_box(slide, 0.9, 0.55, 8.5, 0.8, "O que isso amplia?", size=32, bold=True, color=INK)
    add_text_box(
        slide,
        0.9,
        1.5,
        4.0,
        4.5,
        "Agente para terapeuta\n\nNão é só automação.\nÉ capacidade operacional que antes só estruturas maiores tinham.",
        size=19,
        color=INK,
    )
    add_text_box(
        slide,
        5.2,
        1.5,
        4.0,
        4.5,
        "IA para pequenas empresas\n\nNão é só software.\nÉ distribuir acesso tecnológico e ampliar capacidade produtiva.",
        size=19,
        color=INK,
    )
    shape = slide.shapes.add_shape(1, Inches(4.95), Inches(1.5), Inches(0.04), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = VIOLET
    shape.line.fill.background()
    add_text_box(
        slide,
        0.9,
        5.8,
        8.2,
        0.8,
        "Ferramentas para ampliar capacidades humanas — não substituí-las.",
        size=18,
        bold=True,
        color=VIOLET,
    )
    add_accent_bar(slide)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1 Capa
    slide_title(
        prs,
        "Os Fundamentos da BTECH",
        "Uma conversa sobre o futuro que queremos construir.",
        dark=True,
    )

    # 2 Por quê — ABERTURA
    slide_discussion(
        prs,
        "Por que a BTECH existe?",
        "Antes de tudo: cada fundador responde em voz alta (3 min).\nSem debate ainda — só escutar. Registrar no quadro.",
    )

    # 3 Quem somos
    slide = slide_title(prs, "Quem é a BTECH hoje?", dark=False)
    add_bullets(
        slide,
        0.9,
        3.2,
        8.2,
        3.5,
        [
            "Hoje, a BTECH somos nós.",
            "Não somos um site, infraestrutura, plataforma ou produto.",
            "Somos pessoas tentando compreender uma transformação histórica.",
            "",
            "Objetivo deste encontro:",
            "Descobrir o que estamos realmente construindo.",
        ],
        size=20,
        color=INK,
    )

    # 4 Sonho
    slide_quote(
        prs,
        [
            "Sonho que se sonha só é apenas um sonho.",
            "Sonho que se sonha junto é realidade.",
        ],
    )

    # 5 Timeline
    slide = slide_title(prs, "O mundo em transformação", "Já vivemos algo parecido antes?", dark=False)
    add_text_box(
        slide,
        0.9,
        4.0,
        8.5,
        1.2,
        "Agricultura  →  Revolução Industrial  →  Internet  →  Inteligência Artificial",
        size=24,
        bold=True,
        color=VIOLET,
        align=PP_ALIGN.CENTER,
    )

    # 6 Vapor
    slide = slide_title(prs, "A última grande transformação", dark=False)
    add_bullets(
        slide,
        0.9,
        3.0,
        8.2,
        3.8,
        [
            "A máquina a vapor não mudou apenas fábricas.",
            "Mudou cidades, famílias, educação, profissões, governos, transporte, economia.",
            "O trabalho tornou-se o principal organizador da vida social.",
            "",
            "Não foi só tecnologia — foi reorganização da sociedade.",
        ],
        size=20,
    )

    # 7 E agora
    slide = slide_title(prs, "E agora?", dark=False)
    add_bullets(
        slide,
        0.9,
        2.8,
        8.2,
        4.0,
        [
            "IA · Agentes · Automação · Conhecimento digital",
            "",
            "Conhecimento distribuído instantaneamente.",
            "Software com custo marginal próximo de zero.",
            "Uma pessoa pode ter capacidades antes restritas a grandes organizações.",
            "",
            "Nova tecnologia — ou nova reorganização da sociedade?",
        ],
        size=19,
    )

    # 8 Educação
    slide = slide_title(prs, "Pense na educação", dark=False)
    add_bullets(
        slide,
        0.9,
        2.8,
        8.2,
        4.0,
        [
            "A escola industrial formou pessoas para processos, horários e tarefas repetíveis.",
            "",
            "Uma IA pode ensinar, traduzir, gerar código, analisar dados, tutorar.",
            "",
            "O que significa educar alguém nesse mundo?",
            "Criatividade · pensamento crítico · ética · colaboração · boas perguntas.",
        ],
        size=19,
    )

    # 9 Humano
    slide = slide_title(prs, "Qual é o papel do ser humano?", dark=False)
    add_bullets(
        slide,
        0.9,
        2.8,
        8.2,
        3.8,
        [
            "A IA executa tarefas. O propósito continua sendo humano.",
            "",
            "Talvez a IA não esteja aqui para nos substituir.",
            "Talvez esteja aqui para nos libertar do repetitivo.",
            "",
            "Não substituir pessoas. Potencializar pessoas.",
        ],
        size=20,
    )

    # 10 Dois futuros
    slide_split_futures(prs)

    # 11 Discussão futuro
    slide_discussion(
        prs,
        "Em qual futuro queremos atuar?",
        "15 min · Cada um 2 min + síntese no quadro.",
    )

    # 12 Onde entra
    slide = slide_title(prs, "Onde a BTECH entra?", dark=False)
    add_text_box(
        slide,
        0.9,
        3.0,
        8.2,
        3.0,
        "Cloud? Backup? Infraestrutura? Automação? IA? Agentes?\n\nOu algo maior?",
        size=26,
        color=INK,
    )

    # 13 Construímos
    slide = slide_title(prs, "O que temos construído?", dark=False)
    add_text_box(
        slide,
        0.9,
        3.0,
        8.2,
        1.0,
        "Cloud · Backup · Infraestrutura · Automação · Agentes · IA",
        size=22,
        bold=True,
        color=VIOLET,
    )
    add_text_box(
        slide,
        0.9,
        4.2,
        8.2,
        1.0,
        "O que todas essas iniciativas têm em comum?",
        size=22,
        color=INK,
        italic=True,
    )

    # 14 Exemplos
    slide_examples(prs)

    # 15 Pergunta invertida
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, INK)
    add_text_box(
        slide,
        0.9,
        1.8,
        8.2,
        1.2,
        "A maioria pergunta:",
        size=22,
        color=INK_SOFT,
    )
    add_text_box(
        slide,
        0.9,
        2.5,
        8.2,
        1.5,
        "“Como ganhamos dinheiro com essa tecnologia?”",
        size=26,
        color=WHITE,
        italic=True,
    )
    add_text_box(
        slide,
        0.9,
        4.2,
        8.2,
        1.2,
        "Nossa pergunta:",
        size=22,
        color=CYAN,
        bold=True,
    )
    add_text_box(
        slide,
        0.9,
        4.9,
        8.2,
        1.8,
        "“Como ganhamos dinheiro ampliando capacidades humanas através da tecnologia?”",
        size=26,
        color=WHITE,
        bold=True,
    )
    add_gradient_bar(slide)

    # 16 Negócio
    slide = slide_title(prs, "Importante", dark=False)
    add_bullets(
        slide,
        0.9,
        2.8,
        8.2,
        4.0,
        [
            "A BTECH não é um movimento. É um negócio.",
            "Queremos crescer, gerar valor, lucro e sustentabilidade.",
            "",
            "Acreditamos que crescimento econômico e transformação positiva podem caminhar juntos.",
        ],
        size=21,
    )

    # 17 Significado do B
    slide = slide_title(prs, "O significado do B", dark=True)
    add_text_box(
        slide,
        0.9,
        3.2,
        8.2,
        2.5,
        "O B não é apenas uma alternativa tecnológica.\n\nÉ uma alternativa de visão — uma escolha sobre como construir o futuro.\n\nAmpliação, não concentração.",
        size=24,
        color=INK_SOFT,
    )

    # 18 Discussão meio
    slide_discussion(
        prs,
        "Debate: o que estamos construindo de verdade?",
        "20 min · Reações à proposta abaixo. Ajustar palavras, não só concordar.",
    )

    # 19 Porquê proposta
    slide = slide_title(prs, "Proposta de PORQUÊ", "(rascunho para reação do grupo)", dark=False)
    add_text_box(
        slide,
        0.9,
        3.0,
        8.2,
        3.5,
        "Acreditamos que a próxima revolução tecnológica deve ampliar capacidades humanas, democratizar oportunidades e contribuir para uma sociedade mais autônoma, acessível e próspera.\n\nA BTECH existe para participar ativamente dessa transformação.",
        size=20,
        color=INK,
    )

    # 20 Como + O quê
    slide = slide_title(prs, "Proposta de COMO e O QUÊ", dark=False)
    add_bullets(
        slide,
        0.9,
        2.6,
        8.2,
        4.5,
        [
            "COMO: Transformar tecnologia complexa em solução acessível. Instrumento, não finalidade.",
            "",
            "O QUÊ hoje: Cloud · Backup · Infra · Automação · Agentes",
            "O QUÊ amanhã: Educação · Plataformas · Produtos · Laboratórios · Pesquisa",
            "",
            "O quê muda. O porquê permanece.",
        ],
        size=19,
    )

    # 21 Por quê — FECHAMENTO
    slide_discussion(
        prs,
        "Por que a BTECH existe?",
        "Fechamento · 15 min · Cada um responde de novo, depois do debate.\nO que mudou na sua resposta? Síntese coletiva — 1 frase se possível.",
    )

    # 22 Final
    slide_quote(
        prs,
        [
            "O futuro não é algo que encontramos.",
            "É algo que construímos.",
        ],
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Gerado: {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()

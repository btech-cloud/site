#!/usr/bin/env python3
"""Gera template PPTX 1920×1080 — orçamentos e apresentações BTech."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from presentation_theme import (
    ASSETS,
    BLUE,
    CYAN,
    INK,
    INK_SOFT,
    LIGHT_BG,
    MAGENTA,
    MUTED,
    ORANGE,
    PANEL,
    ROOT,
    VIOLET,
    WHITE,
    add_bg_image,
    add_bullets,
    add_flow_nodes,
    add_footer,
    add_gradient_rule,
    add_hex_accent,
    add_logo,
    add_multiline,
    add_number_highlight,
    add_photo_panel,
    add_slide_label,
    add_table,
    add_text_box,
    generate_network_bg,
    set_slide_bg,
    set_slide_size,
)

OUT = ROOT / "assets" / "presentations" / "btech-template-orcamento.pptx"

# Imagens corporativas / pessoas
IMG_TEAM = ROOT / "assets" / "btech-operations-team.png"
IMG_WORKSHOP = ROOT / "assets" / "btech-people-workshop.png"
IMG_CONVERSA = ROOT / "assets" / "btech-client-conversation.png"
IMG_MEETING = ROOT / "assets" / "images" / "equipe-trabalho.jpg"


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def slide_cover(prs, bg_dark: Path) -> None:
    slide = blank(prs)
    add_bg_image(slide, bg_dark)
    add_logo(slide, light_bg=False)
    add_text_box(slide, 0.7, 2.0, 8.5, 1.2, "Proposta comercial", size=14, color=CYAN, bold=True)
    add_text_box(
        slide,
        0.7,
        2.5,
        10.5,
        1.8,
        "Amplie sua capacidade\nna era da Inteligência Artificial",
        size=40,
        bold=True,
        color=WHITE,
    )
    add_gradient_rule(slide, left=0.7, top=4.35, width=4.0)
    add_text_box(
        slide,
        0.7,
        4.7,
        9.0,
        1.0,
        "Cloud · Backup · Automação · Agentes de IA",
        size=18,
        color=INK_SOFT,
    )
    add_text_box(slide, 0.7, 5.8, 8.0, 0.5, "Cliente: [Nome da empresa]", size=16, color=MUTED)
    add_text_box(slide, 0.7, 6.2, 8.0, 0.5, "Ref.: [BT-2026-001]  ·  [Mês/Ano]", size=14, color=MUTED)
    add_footer(slide, "BTech · Confidencial · Uso exclusivo do destinatário")


def slide_agenda(prs, bg_light: Path) -> None:
    slide = blank(prs)
    set_slide_bg(slide, LIGHT_BG)
    add_bg_image(slide, bg_light)
    add_logo(slide, light_bg=True)
    add_slide_label(slide, "Sumário", light=True)
    add_text_box(slide, 0.7, 1.0, 10, 0.9, "O que vamos apresentar", size=34, bold=True, color=INK)
    add_gradient_rule(slide, left=0.7, top=1.85, width=2.8)
    items = [
        "01  Contexto e desafio",
        "02  Nossa proposta de valor",
        "03  Solução e arquitetura",
        "04  Escopo, fases e cronograma",
        "05  Investimento e condições",
        "06  Por que a BTech",
        "07  Próximos passos",
    ]
    add_bullets(slide, 0.9, 2.3, 6.5, 4.5, items, size=20, color=INK_SOFT)
    add_photo_panel(slide, IMG_CONVERSA, left=7.2, top=1.8, width=5.5, height=4.6)
    add_footer(slide, light=True)


def slide_about(prs, bg_light: Path) -> None:
    slide = blank(prs)
    set_slide_bg(slide, LIGHT_BG)
    add_bg_image(slide, bg_light)
    add_logo(slide, light_bg=True)
    add_slide_label(slide, "Quem somos", light=True)
    add_text_box(slide, 0.7, 1.0, 6.2, 1.0, "Tecnologia humana,\nresultado real", size=32, bold=True, color=INK)
    add_multiline(
        slide,
        0.7,
        2.2,
        6.0,
        4.5,
        [
            (
                "A BTech é uma empresa de tecnologia, infraestrutura e inteligência artificial focada em ampliar capacidades humanas e organizacionais.",
                {"size": 17, "color": INK_SOFT},
            ),
            ("", {}),
            (
                "Construímos infraestrutura, automação e agentes inteligentes — com consultoria próxima da operação, do diagnóstico ao go-live.",
                {"size": 17, "color": INK_SOFT},
            ),
            ("", {}),
            ("Não somos provedor de datacenter. Somos o parceiro que organiza, desenha e acompanha.", {"size": 16, "bold": True, "color": VIOLET}),
        ],
    )
    add_photo_panel(slide, IMG_TEAM, left=7.0, top=1.4, width=5.8, height=5.0)
    add_footer(slide, light=True)


def slide_why_b(prs, bg_dark: Path) -> None:
    slide = blank(prs)
    add_bg_image(slide, bg_dark)
    add_logo(slide)
    add_slide_label(slide, "Por que B")
    add_text_box(
        slide,
        0.7,
        1.0,
        11,
        1.0,
        "O B da BTech — alternativa com propósito",
        size=32,
        bold=True,
        color=WHITE,
    )
    add_gradient_rule(slide, left=0.7, top=1.85, width=3.0)
    cards = [
        ("Bridge", "Pontes entre negócio, TI e operação"),
        ("Build", "Construímos soluções, não só slides"),
        ("Beta", "Iteração rápida com governança"),
        ("Brasil", "Proximidade, contexto e confiança"),
        ("Better", "Ampliação, não concentração"),
        ("Backup", "Continuidade como base, não opcional"),
    ]
    x0, y0 = 0.7, 2.3
    for i, (title, desc) in enumerate(cards):
        col = i % 3
        row = i // 3
        x = x0 + col * 4.0
        y = y0 + row * 2.1
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.6), Inches(1.75))
        card.fill.solid()
        card.fill.fore_color.rgb = PANEL
        card.line.color.rgb = VIOLET
        card.line.width = Pt(1)
        add_text_box(slide, x + 0.2, y + 0.2, 3.2, 0.5, title, size=22, bold=True, color=CYAN)
        add_text_box(slide, x + 0.2, y + 0.75, 3.2, 0.9, desc, size=14, color=INK_SOFT)
    add_footer(slide)


def slide_context(prs, bg_light: Path) -> None:
    slide = blank(prs)
    set_slide_bg(slide, LIGHT_BG)
    add_bg_image(slide, bg_light)
    add_logo(slide, light_bg=True)
    add_slide_label(slide, "Contexto", light=True)
    add_text_box(slide, 0.7, 1.0, 10, 0.8, "Entendimento do cenário", size=34, bold=True, color=INK)
    add_text_box(
        slide,
        0.7,
        1.9,
        11.5,
        0.8,
        "[Cliente] opera em um ambiente onde cloud, dados e IA deixaram de ser opcionais — mas a complexidade operacional cresceu mais rápido que a capacidade interna.",
        size=18,
        color=INK_SOFT,
    )
    add_number_highlight(slide, 0.7, 3.2, "73%", "processos ainda manuais*", accent=CYAN)
    add_number_highlight(slide, 3.5, 3.2, "4+", "sistemas sem integração*", accent=VIOLET)
    add_number_highlight(slide, 6.3, 3.2, "∞", "pressão por continuidade*", accent=ORANGE)
    add_text_box(slide, 0.7, 5.1, 11, 0.4, "* Exemplo — substituir por dados reais do diagnóstico", size=10, italic=True, color=MUTED)
    add_bullets(
        slide,
        0.7,
        5.5,
        11,
        1.5,
        [
            "Gestão precisa de visibilidade sem depender só da TI.",
            "TI precisa de arquitetura sustentável e parceiros que co-entreguem.",
            "Operação não pode parar enquanto a transformação acontece.",
        ],
        size=16,
        color=INK_SOFT,
    )
    add_footer(slide, light=True)


def slide_challenges(prs) -> None:
    slide = blank(prs)
    set_slide_bg(slide, PANEL)
    add_slide_label(slide, "Desafios")
    add_text_box(slide, 0.7, 1.0, 10, 0.8, "O que precisa mudar", size=34, bold=True, color=WHITE)
    challenges = [
        ("01", "Infraestrutura fragmentada", "Ambientes críticos sem padrão, documentação ou dono claro."),
        ("02", "Risco de continuidade", "Backup e DR não testados com a frequência que o negócio exige."),
        ("03", "Tarefas repetitivas", "Equipe absorve trabalho que deveria estar automatizado."),
        ("04", "IA sem governança", "Experimentos isolados, sem integração nem supervisão humana."),
    ]
    y = 2.2
    for num, title, desc in challenges:
        add_text_box(slide, 0.7, y, 0.6, 0.5, num, size=20, bold=True, color=CYAN)
        add_text_box(slide, 1.4, y, 4.5, 0.4, title, size=18, bold=True, color=WHITE)
        add_text_box(slide, 1.4, y + 0.42, 10, 0.5, desc, size=14, color=INK_SOFT)
        y += 1.15
    add_hex_accent(slide, 4, side="right")
    add_footer(slide)


def slide_objectives(prs, bg_light: Path) -> None:
    slide = blank(prs)
    set_slide_bg(slide, LIGHT_BG)
    add_bg_image(slide, bg_light)
    add_logo(slide, light_bg=True)
    add_slide_label(slide, "Objetivos", light=True)
    add_text_box(slide, 0.7, 1.0, 10, 0.8, "O que este projeto precisa alcançar", size=32, bold=True, color=INK)
    add_bullets(
        slide,
        0.7,
        2.0,
        6.0,
        4.5,
        [
            "Estabilizar e modernizar a base cloud com governança.",
            "Garantir proteção e recuperação testada dos dados críticos.",
            "Automatizar fluxos repetitivos com integração entre sistemas.",
            "Introduzir agentes de IA com supervisão e métricas claras.",
            "Capacitar a equipe para sustentar a operação após o go-live.",
        ],
        size=17,
        color=INK_SOFT,
    )
    add_photo_panel(slide, IMG_WORKSHOP, left=7.0, top=1.6, width=5.8, height=4.8)
    add_footer(slide, light=True)


def slide_solution_vision(prs, bg_dark: Path) -> None:
    slide = blank(prs)
    add_bg_image(slide, bg_dark)
    add_logo(slide)
    add_slide_label(slide, "Proposta")
    add_text_box(slide, 0.7, 1.8, 11, 1.5, "Uma solução que amplia capacidade\n— não adiciona complexidade", size=36, bold=True, color=WHITE)
    add_text_box(
        slide,
        0.7,
        3.6,
        10,
        1.2,
        "Combinamos infraestrutura confiável, automação pragmática e IA supervisionada, entregues em fases que a operação consegue absorver.",
        size=20,
        color=INK_SOFT,
    )
    add_gradient_rule(slide, left=0.7, top=5.0, width=4.5)
    add_footer(slide)


def slide_architecture_flow(prs, bg_light: Path) -> None:
    slide = blank(prs)
    set_slide_bg(slide, LIGHT_BG)
    add_bg_image(slide, bg_light)
    add_logo(slide, light_bg=True)
    add_slide_label(slide, "Arquitetura", light=True)
    add_text_box(slide, 0.7, 1.0, 10, 0.7, "Fluxo da solução proposta", size=32, bold=True, color=INK)
    nodes = [
        {"step": "01", "label": "Usuários\n& Gestão", "fill": PANEL, "stroke": VIOLET, "weight": 3},
        {"step": "02", "label": "Automação\n& Integração", "fill": PANEL, "stroke": BLUE, "weight": 2.5},
        {"step": "03", "label": "Cloud\n& Dados", "fill": PANEL, "stroke": CYAN, "weight": 3.5},
        {"step": "04", "label": "Agentes\nIA", "fill": PANEL, "stroke": MAGENTA, "weight": 2},
        {"step": "05", "label": "Backup\n& DR", "fill": PANEL, "stroke": ORANGE, "weight": 2.5},
    ]
    add_flow_nodes(slide, nodes, start_x=0.65, y=2.8)
    add_text_box(
        slide,
        0.7,
        4.5,
        12,
        1.5,
        "Linhas mais espessas = integrações críticas  ·  Hexágonos = pontos de conexão na operação\nSubstitua os rótulos pelos sistemas reais do cliente (ERP, CRM, cloud provider, etc.)",
        size=13,
        color=MUTED,
        italic=True,
    )
    add_footer(slide, light=True)


def slide_scope_table(prs) -> None:
    slide = blank(prs)
    set_slide_bg(slide, LIGHT_BG)
    add_logo(slide, light_bg=True)
    add_slide_label(slide, "Escopo", light=True)
    add_text_box(slide, 0.7, 1.0, 10, 0.7, "Escopo de trabalho", size=32, bold=True, color=INK)
    data = [
        ["Frente", "Entrega", "Responsável"],
        ["Cloud", "Arquitetura, migração e runbook operacional", "BTech + [Parceiro cloud]"],
        ["Backup / DR", "Política, implementação e teste de restore", "BTech"],
        ["Automação", "Fluxos n8n / CI-CD / integrações API", "BTech"],
        ["Agentes IA", "Agente atendimento + base de conhecimento", "BTech"],
        ["Governança", "Rituais, indicadores e documentação", "BTech + Cliente"],
    ]
    add_table(slide, 0.7, 2.0, 11.8, len(data), 3, data)
    add_footer(slide, light=True)


def slide_phases_timeline(prs, bg_dark: Path) -> None:
    slide = blank(prs)
    add_bg_image(slide, bg_dark)
    add_logo(slide)
    add_slide_label(slide, "Cronograma")
    add_text_box(slide, 0.7, 1.0, 10, 0.7, "Fases e marcos", size=32, bold=True, color=WHITE)
    phases = [
        ("Fase 1", "Diagnóstico", "Semanas 1–2", "Mapa de sistemas, riscos e quick wins"),
        ("Fase 2", "Fundação", "Semanas 3–6", "Cloud, backup e primeiras automações"),
        ("Fase 3", "Escala", "Semanas 7–10", "Agentes IA e integrações avançadas"),
        ("Fase 4", "Operação", "Semanas 11–12", "Go-live, handover e hypercare"),
    ]
    y = 2.1
    for tag, name, when, desc in phases:
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(y), Inches(11.8), Inches(0.95))
        bar.fill.solid()
        bar.fill.fore_color.rgb = PANEL
        bar.line.color.rgb = VIOLET
        add_text_box(slide, 0.9, y + 0.12, 1.2, 0.4, tag, size=11, bold=True, color=CYAN)
        add_text_box(slide, 2.1, y + 0.1, 2.5, 0.4, name, size=16, bold=True, color=WHITE)
        add_text_box(slide, 4.8, y + 0.12, 2.0, 0.4, when, size=12, color=ORANGE)
        add_text_box(slide, 7.0, y + 0.12, 5.0, 0.5, desc, size=13, color=INK_SOFT)
        y += 1.15
    add_footer(slide)


def slide_deliverables(prs) -> None:
    slide = blank(prs)
    set_slide_bg(slide, LIGHT_BG)
    add_logo(slide, light_bg=True)
    add_slide_label(slide, "Entregáveis", light=True)
    add_text_box(slide, 0.7, 1.0, 10, 0.7, "O que você recebe ao final de cada fase", size=30, bold=True, color=INK)
    data = [
        ["Entregável", "Formato", "Fase"],
        ["Documento de arquitetura (HLD/LLD)", "PDF + diagrama", "1–2"],
        ["Ambiente cloud provisionado", "Console + IaC", "2"],
        ["Política e job de backup testado", "Relatório de restore", "2"],
        ["Fluxos de automação documentados", "Repositório + runbook", "2–3"],
        ["Agente IA em produção supervisionada", "Painel + métricas", "3"],
        ["Handover e plano de operação", "Workshop + checklist", "4"],
    ]
    add_table(slide, 0.7, 1.9, 11.8, len(data), 3, data)
    add_footer(slide, light=True)


def slide_investment(prs, bg_dark: Path) -> None:
    slide = blank(prs)
    add_bg_image(slide, bg_dark)
    add_logo(slide)
    add_slide_label(slide, "Investimento")
    add_text_box(slide, 0.7, 1.0, 10, 0.7, "Proposta econômica", size=34, bold=True, color=WHITE)
    data = [
        ["Item", "Descrição", "Valor (R$)"],
        ["Fase 1 — Diagnóstico", "Workshops, mapeamento e plano executivo", "[ 28.000 ]"],
        ["Fase 2 — Fundação", "Cloud, backup e automações base", "[ 85.000 ]"],
        ["Fase 3 — Escala", "Agentes IA e integrações", "[ 62.000 ]"],
        ["Fase 4 — Operação", "Go-live, hypercare e handover", "[ 24.000 ]"],
        ["Total do projeto", "12 semanas · escopo conforme anexo", "[ 199.000 ]"],
    ]
    add_table(slide, 0.7, 1.9, 11.8, len(data), 3, data)
    add_text_box(slide, 0.7, 5.0, 6, 0.5, "Investimento mensal recorrente (opcional):", size=14, bold=True, color=CYAN)
    add_text_box(slide, 0.7, 5.45, 8, 0.5, "Suporte e evolução contínua — [ R$ 8.500 / mês ]", size=16, color=WHITE)
    add_text_box(slide, 0.7, 6.1, 11, 0.5, "Valores de referência · ajustar conforme escopo real · impostos conforme regime BTech", size=10, italic=True, color=MUTED)
    add_footer(slide)


def slide_conditions(prs) -> None:
    slide = blank(prs)
    set_slide_bg(slide, LIGHT_BG)
    add_logo(slide, light_bg=True)
    add_slide_label(slide, "Condições", light=True)
    add_text_box(slide, 0.7, 1.0, 10, 0.7, "Premissas e condições comerciais", size=30, bold=True, color=INK)
    add_bullets(
        slide,
        0.7,
        2.0,
        5.5,
        4.8,
        [
            "Validade da proposta: 30 dias",
            "Pagamento: 30% na assinatura, 40% na Fase 2, 30% na conclusão",
            "Prazos dependem de acesso a ambientes e stakeholders",
            "Licenças de cloud e software são do cliente",
            "Horas adicionais: [ R$ / h ] sob aprovação prévia",
        ],
        title="Condições",
        size=15,
        color=INK_SOFT,
    )
    add_bullets(
        slide,
        6.8,
        2.0,
        5.5,
        4.8,
        [
            "Fora de escopo: desenvolvimento de software sob medida não previsto",
            "Exclusões: hardware físico e links de telecom",
            "Ambientes legados sem documentação podem exigir fase extra",
            "Dados sensíveis seguem política LGPD acordada",
            "Confidencialidade mútua em anexo",
        ],
        title="Premissas e exclusões",
        size=15,
        color=INK_SOFT,
    )
    add_footer(slide, light=True)


def slide_why_btech(prs, bg_light: Path) -> None:
    slide = blank(prs)
    set_slide_bg(slide, LIGHT_BG)
    add_bg_image(slide, bg_light)
    add_logo(slide, light_bg=True)
    add_slide_label(slide, "Diferenciais", light=True)
    add_text_box(slide, 0.7, 1.0, 10, 0.7, "Por que a BTech", size=32, bold=True, color=INK)
    add_number_highlight(slide, 0.7, 2.2, "1", "Conversa antes da ferramenta", accent=LIGHT_BG)
    add_number_highlight(slide, 3.5, 2.2, "2", "Grupo consultivo multidisciplinar", accent=LIGHT_BG)
    add_number_highlight(slide, 6.3, 2.2, "3", "Co-entrega com parceiros", accent=LIGHT_BG)
    add_number_highlight(slide, 9.1, 2.2, "4", "IA com supervisão humana", accent=LIGHT_BG)
    add_multiline(
        slide,
        0.7,
        4.2,
        11.5,
        2.5,
        [
            ("Não vendemos datacenter. Vendemos resultado.", {"size": 20, "bold": True, "color": VIOLET}),
            ("", {}),
            (
                "Experiência em AWS, Azure, GCP, Kubernetes, OpenShift, Veeam, SAS e automação — com linguagem que gestão e TI entendem juntos.",
                {"size": 16, "color": INK_SOFT},
            ),
        ],
    )
    add_photo_panel(slide, IMG_MEETING, left=8.5, top=4.0, width=4.2, height=2.8)
    add_footer(slide, light=True)


def slide_next_steps(prs, bg_dark: Path) -> None:
    slide = blank(prs)
    add_bg_image(slide, bg_dark)
    add_logo(slide)
    add_slide_label(slide, "Próximos passos")
    add_text_box(slide, 0.7, 1.0, 10, 0.7, "Como seguimos", size=34, bold=True, color=WHITE)
    steps = [
        {"step": "1", "label": "Alinhar\nescopo"},
        {"step": "2", "label": "Ajustar\ninvestimento"},
        {"step": "3", "label": "Assinar\ncontrato"},
        {"step": "4", "label": "Kick-off\nFase 1"},
    ]
    add_flow_nodes(slide, steps, start_x=1.0, y=2.6)
    add_text_box(slide, 0.7, 4.5, 11, 0.8, "Próxima reunião sugerida: [data] · Participantes: [nomes]", size=16, color=INK_SOFT)
    add_footer(slide)


def slide_closing(prs, bg_dark: Path) -> None:
    slide = blank(prs)
    add_bg_image(slide, bg_dark)
    add_logo(slide, stacked=True)
    add_text_box(
        slide,
        0.7,
        2.4,
        11,
        1.5,
        "O futuro não é algo que encontramos.\nÉ algo que construímos.",
        size=36,
        bold=True,
        color=WHITE,
    )
    add_gradient_rule(slide, left=0.7, top=4.2, width=4.0)
    add_multiline(
        slide,
        0.7,
        4.6,
        8,
        2,
        [
            ("comercial@b-tech.cloud", {"size": 18, "color": CYAN}),
            ("(11) 9 3022-6495", {"size": 16, "color": INK_SOFT}),
            ("btech.cloud · Brasília/DF", {"size": 14, "color": MUTED}),
        ],
    )
    add_footer(slide, "BTech · Obrigado pela confiança")


def slide_layout_guide(prs) -> None:
    """Slide guia interno — tipos de layout disponíveis."""
    slide = blank(prs)
    set_slide_bg(slide, PANEL)
    add_text_box(slide, 0.7, 0.5, 12, 0.6, "GUIA DO TEMPLATE (excluir antes de enviar ao cliente)", size=14, bold=True, color=ORANGE)
    add_text_box(slide, 0.7, 1.2, 12, 5.5, """Layouts incluídos neste arquivo:

• Capa escura — hex + teia + logo
• Sumário / agenda
• Quem somos — texto + foto corporativa
• Por que B — 6 cards (Bridge, Build, Beta…)
• Contexto — destaque numérico + bullets
• Desafios — lista numerada fundo escuro
• Objetivos — bullets + foto workshop
• Visão da solução — statement escuro
• Arquitetura — fluxo hexagonal conectado
• Escopo — tabela editável
• Cronograma — fases em faixas
• Entregáveis — tabela
• Investimento — tabela + recorrência
• Condições — premissas duas colunas
• Por que BTech — 4 diferenciais + foto
• Próximos passos — fluxo 4 etapas
• Encerramento — CTA + contato

Formato: 1920×1080 (16:9) · Fonte sugerida: Calibri
Fundos gerados: assets/presentations/template/bg-*.png
Edite no PowerPoint, Google Slides ou importe no Canva.""", size=13, color=INK_SOFT)
    add_footer(slide)


def build() -> None:
    bg_dark = generate_network_bg("bg-dark-network.png", dark=True, seed=42)
    bg_light = generate_network_bg("bg-light-network.png", dark=False, seed=77)

    prs = Presentation()
    set_slide_size(prs)

    slide_layout_guide(prs)
    slide_cover(prs, bg_dark)
    slide_agenda(prs, bg_light)
    slide_about(prs, bg_light)
    slide_why_b(prs, bg_dark)
    slide_context(prs, bg_light)
    slide_challenges(prs)
    slide_objectives(prs, bg_light)
    slide_solution_vision(prs, bg_dark)
    slide_architecture_flow(prs, bg_light)
    slide_scope_table(prs)
    slide_phases_timeline(prs, bg_dark)
    slide_deliverables(prs)
    slide_investment(prs, bg_dark)
    slide_conditions(prs)
    slide_why_btech(prs, bg_light)
    slide_next_steps(prs, bg_dark)
    slide_closing(prs, bg_dark)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Gerado: {OUT}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Fundos: {bg_dark.name}, {bg_light.name}")


if __name__ == "__main__":
    build()

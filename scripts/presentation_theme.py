"""Tema visual BTech para apresentações 1920×1080 — hex, teia, gradientes."""

from __future__ import annotations

import math
import random
from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "assets" / "presentations" / "template"
ASSETS.mkdir(parents=True, exist_ok=True)

# 16:9 @ 1920×1080
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
PX_W, PX_H = 1920, 1080

# Paleta Instagram / site
INK = RGBColor(5, 5, 16)
INK_SOFT = RGBColor(36, 50, 79)
PANEL = RGBColor(26, 11, 64)
PANEL_DEEP = RGBColor(15, 8, 36)
VIOLET = RGBColor(110, 30, 230)
BLUE = RGBColor(20, 60, 255)
CYAN = RGBColor(0, 225, 180)
TEAL = RGBColor(0, 225, 255)
MAGENTA = RGBColor(200, 0, 210)
ORANGE = RGBColor(255, 175, 60)
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(250, 251, 255)
LILAC = RGBColor(243, 240, 255)
MUTED = RGBColor(120, 115, 145)
SUCCESS = RGBColor(90, 232, 56)

PALETTE_RGB = [
    (20, 60, 255),
    (110, 30, 230),
    (200, 0, 210),
    (0, 225, 255),
    (0, 225, 180),
    (255, 175, 60),
]


def hex_points(cx: float, cy: float, r: float) -> list[tuple[float, float]]:
    return [
        (
            cx + r * math.cos(math.pi / 3 * i - math.pi / 6),
            cy + r * math.sin(math.pi / 3 * i - math.pi / 6),
        )
        for i in range(6)
    ]


def _lerp(a: int, b: int, t: float) -> int:
    return int(a + (b - a) * t)


def _gradient_color(t: float) -> tuple[int, int, int, int]:
    """Gradiente violeta → azul → ciano."""
    stops = [(0.0, PALETTE_RGB[1]), (0.45, PALETTE_RGB[0]), (1.0, PALETTE_RGB[4])]
    if t <= stops[0][0]:
        c = stops[0][1]
    elif t >= stops[-1][0]:
        c = stops[-1][1]
    else:
        for i in range(len(stops) - 1):
            t0, c0 = stops[i]
            t1, c1 = stops[i + 1]
            if t0 <= t <= t1:
                u = (t - t0) / (t1 - t0)
                c = tuple(_lerp(c0[j], c1[j], u) for j in range(3))
                break
        else:
            c = stops[-1][1]
    return (*c, 255)


def generate_network_bg(
    name: str,
    *,
    dark: bool = True,
    seed: int = 42,
) -> Path:
    """Gera fundo com hexágonos, brilhos e teia de conexões."""
    out = ASSETS / name
    if out.exists():
        return out

    rng = random.Random(seed)
    base = (5, 5, 16) if dark else (250, 251, 255)
    img = Image.new("RGBA", (PX_W, PX_H), (*base, 255))
    draw = ImageDraw.Draw(img)

    # Orbs de luz (bokeh)
    for _ in range(8 if dark else 5):
        cx = rng.randint(0, PX_W)
        cy = rng.randint(0, PX_H)
        r = rng.randint(120, 420)
        color = PALETTE_RGB[rng.randint(0, len(PALETTE_RGB) - 1)]
        alpha = rng.randint(18, 55) if dark else rng.randint(8, 22)
        orb = Image.new("RGBA", (r * 2, r * 2), (0, 0, 0, 0))
        od = ImageDraw.Draw(orb)
        od.ellipse((0, 0, r * 2, r * 2), fill=(*color, alpha))
        orb = orb.filter(ImageFilter.GaussianBlur(radius=r // 3))
        img.alpha_composite(orb, (cx - r, cy - r))

    # Nós da teia
    nodes: list[tuple[float, float, float]] = []
    for _ in range(22):
        nodes.append((rng.uniform(80, PX_W - 80), rng.uniform(80, PX_H - 80), rng.uniform(0.3, 1.0)))

    # Conexões com pesos diferentes
    draw = ImageDraw.Draw(img)
    for i, (x1, y1, w1) in enumerate(nodes):
        for j, (x2, y2, w2) in enumerate(nodes):
            if j <= i:
                continue
            dist = math.hypot(x2 - x1, y2 - y1)
            if dist > 380 or rng.random() > 0.12:
                continue
            strength = (w1 + w2) / 2
            alpha = int(35 + 90 * strength) if dark else int(20 + 50 * strength)
            width = max(1, int(1 + 3 * strength))
            t = (x1 + x2) / (2 * PX_W)
            col = _gradient_color(t)[:3]
            draw.line((x1, y1, x2, y2), fill=(*col, alpha), width=width)

    # Hexágonos — cheios, desfocados e sutis
    hex_layer = Image.new("RGBA", (PX_W, PX_H), (0, 0, 0, 0))
    hd = ImageDraw.Draw(hex_layer)
    for _ in range(28):
        cx = rng.uniform(0, PX_W)
        cy = rng.uniform(0, PX_H)
        r = rng.uniform(18, 95)
        t = cx / PX_W
        col = _gradient_color(t)
        kind = rng.random()
        if kind < 0.25:
            # Brilho desfocado
            alpha = rng.randint(40, 110)
            pts = hex_points(cx, cy, r)
            hd.polygon(pts, fill=(*col[:3], alpha))
        elif kind < 0.55:
            # Contorno fino
            alpha = rng.randint(50, 140) if dark else rng.randint(30, 90)
            pts = hex_points(cx, cy, r * 0.7)
            hd.polygon(pts, outline=(*col[:3], alpha), width=max(1, int(r / 25)))
        else:
            # Preenchimento sutil
            alpha = rng.randint(25, 75) if dark else rng.randint(15, 45)
            pts = hex_points(cx, cy, r * 0.55)
            hd.polygon(pts, fill=(*col[:3], alpha))

    blur_r = 6 if dark else 4
    hex_layer = hex_layer.filter(ImageFilter.GaussianBlur(radius=blur_r))
    img = Image.alpha_composite(img, hex_layer)

    # Vinheta suave
    vignette = Image.new("L", (PX_W, PX_H), 0)
    vd = ImageDraw.Draw(vignette)
    vd.ellipse((-200, -100, PX_W + 200, PX_H + 200), fill=200)
    vignette = vignette.filter(ImageFilter.GaussianBlur(120))
    if dark:
        dark_layer = Image.new("RGBA", (PX_W, PX_H), (0, 0, 0, 0))
        for y in range(0, PX_H, 4):
            for x in range(0, PX_W, 4):
                v = 255 - vignette.getpixel((x, y))
                if v > 8:
                    dark_layer.putpixel((x, y), (0, 0, 0, min(80, v // 3)))
        img = Image.alpha_composite(img, dark_layer)

    img.convert("RGB").save(out, "PNG", optimize=True)
    return out


def set_slide_size(prs) -> None:
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H


def set_slide_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_bg_image(slide, path: Path) -> None:
    slide.shapes.add_picture(str(path), 0, 0, width=SLIDE_W, height=SLIDE_H)


def add_gradient_rule(slide, left=1.0, top=6.72, width=3.2) -> None:
    colors = [VIOLET, BLUE, CYAN, ORANGE]
    w = width / len(colors)
    x = left
    for c in colors:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(top), Inches(w), Inches(0.07))
        bar.fill.solid()
        bar.fill.fore_color.rgb = c
        bar.line.fill.background()
        x += w


def add_footer(slide, text: str = "BTech · Proposta comercial · Confidencial", *, light: bool = False) -> None:
    color = MUTED if light else RGBColor(160, 155, 185)
    box = slide.shapes.add_textbox(Inches(0.7), Inches(7.05), Inches(8), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = color


def add_slide_label(slide, text: str, *, light: bool = False) -> None:
    color = VIOLET if light else CYAN
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(5), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Calibri"


def add_logo(slide, *, light_bg: bool = False, stacked: bool = False) -> None:
    if stacked:
        path = ROOT / "assets" / "logo-btech-stacked-oficial-light.png" if light_bg else ROOT / "assets" / "logo-btech-stacked-oficial.png"
    else:
        path = ROOT / "assets" / "logo-btech-horizontal-oficial-light.png" if light_bg else ROOT / "assets" / "logo-btech-horizontal-oficial.png"
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(0.65), Inches(0.35), height=Inches(0.42))


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size=24,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    italic=False,
    font_name="Calibri",
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf


def add_multiline(
    slide,
    left,
    top,
    width,
    height,
    lines: list[tuple[str, dict]],
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for text, style in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.font.size = Pt(style.get("size", 20))
        p.font.bold = style.get("bold", False)
        p.font.italic = style.get("italic", False)
        p.font.color.rgb = style.get("color", INK)
        p.font.name = style.get("font", "Calibri")
        p.space_after = Pt(style.get("space_after", 6))
        if "align" in style:
            p.alignment = style["align"]
    return tf


def add_bullets(slide, left, top, width, height, items, *, size=18, color=INK, title=None, bullet_color=VIOLET):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(size + 4)
        p.font.bold = True
        p.font.color.rgb = color
        first = False
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if not item:
            p.text = ""
            continue
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
    return tf


def add_number_highlight(slide, left, top, number: str, label: str, *, accent=CYAN) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(2.6), Inches(1.55))
    card.fill.solid()
    card.fill.fore_color.rgb = PANEL if accent == CYAN else LIGHT_BG
    card.line.color.rgb = accent
    card.line.width = Pt(1.5)
    add_text_box(slide, left + 0.2, top + 0.15, 2.2, 0.8, number, size=36, bold=True, color=accent if accent != LIGHT_BG else VIOLET, align=PP_ALIGN.CENTER)
    add_text_box(slide, left + 0.15, top + 0.95, 2.3, 0.5, label, size=11, color=MUTED if accent == LIGHT_BG else INK_SOFT, align=PP_ALIGN.CENTER)


def add_flow_nodes(slide, nodes: list[dict], *, start_x=1.0, y=3.2) -> None:
    """Fluxo horizontal com hexágonos conectados."""
    n = len(nodes)
    gap = 8.8 / max(n, 1)
    prev_x = None
    for i, node in enumerate(nodes):
        x = start_x + i * gap
        hex_shape = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(x), Inches(y), Inches(1.35), Inches(1.15))
        hex_shape.fill.solid()
        hex_shape.fill.fore_color.rgb = node.get("fill", PANEL)
        hex_shape.line.color.rgb = node.get("stroke", VIOLET)
        hex_shape.line.width = Pt(2)
        add_text_box(slide, x + 0.05, y + 0.22, 1.25, 0.7, node["label"], size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if node.get("step"):
            add_text_box(slide, x + 0.45, y - 0.35, 0.5, 0.3, node["step"], size=10, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
        if prev_x is not None:
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(prev_x + 1.35),
                Inches(y + 0.55),
                Inches(x),
                Inches(y + 0.55),
            )
            conn.line.color.rgb = node.get("line", BLUE)
            conn.line.width = Pt(node.get("weight", 2.5))
        prev_x = x


def add_table(slide, left, top, width, rows, cols, data, *, header=True) -> None:
    height = 0.38 + 0.34 * rows
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11 if r else 12)
                p.font.name = "Calibri"
                p.font.bold = bool(header and r == 0)
                if r == 0:
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = INK if c == 0 else INK_SOFT
            if header and r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PANEL
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LILAC
    return table


def add_photo_panel(slide, image_path: Path, *, left=7.0, top=1.5, width=5.8, height=4.8) -> None:
    if not image_path.exists():
        return
    pic = slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))
    ratio = pic.height / pic.width
    pic.height = Inches(height)
    pic.width = int(pic.height / ratio)
    # Máscara arredondada simulada com overlay de borda
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left - 0.05), Inches(top - 0.05), Inches(width + 0.1), Inches(height + 0.1))
    frame.fill.background()
    frame.line.color.rgb = VIOLET
    frame.line.width = Pt(1.25)


def add_hex_accent(slide, count: int = 3, *, side: str = "right") -> None:
    """Hexágonos decorativos sutis no canto."""
    base_x = 11.2 if side == "right" else 0.4
    for i in range(count):
        y = 0.8 + i * 0.55
        size = 0.35 - i * 0.04
        h = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(base_x), Inches(y), Inches(size), Inches(size * 0.86))
        h.fill.solid()
        colors = [VIOLET, BLUE, CYAN, MAGENTA]
        h.fill.fore_color.rgb = colors[i % len(colors)]
        h.fill.transparency = 0.35 + i * 0.12
        h.line.fill.background()

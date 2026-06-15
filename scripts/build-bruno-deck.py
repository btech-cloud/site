#!/usr/bin/env python3
"""PPTX curto — validação do conceito com Bruno (sem ritual de grupo)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "assets" / "presentations" / "btech-fundamentos-bruno.pptx"

INK = RGBColor(5, 5, 16)
INK_SOFT = RGBColor(244, 242, 255)
PANEL = RGBColor(26, 11, 64)
VIOLET = RGBColor(110, 30, 230)
BLUE = RGBColor(20, 60, 255)
CYAN = RGBColor(0, 225, 180)
ORANGE = RGBColor(255, 175, 60)
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(250, 251, 255)
MUTED = RGBColor(120, 115, 145)


def set_bg(slide, rgb):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = rgb


def bar(slide):
    x = 1.2
    for c in [VIOLET, BLUE, CYAN, ORANGE]:
        s = slide.shapes.add_shape(1, Inches(x), Inches(6.78), Inches(1.55), Inches(0.1))
        s.fill.solid()
        s.fill.fore_color.rgb = c
        s.line.fill.background()
        x += 1.6


def tb(slide, l, t, w, h, text, *, size=24, bold=False, color=INK, center=False, italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT


def bullets(slide, l, t, w, h, items, size=19, color=INK):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)


def title_slide(prs, title, sub=None, dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, INK if dark else LIGHT_BG)
    tc, sc = (WHITE, INK_SOFT) if dark else (INK, MUTED)
    tb(s, 0.85, 2.0, 8.4, 1.2, title, size=38, bold=True, color=tc)
    if sub:
        tb(s, 0.85, 3.1, 8.2, 1.8, sub, size=20, color=sc)
    if dark:
        bar(s)
    return s


def question_slide(prs, title, hint, tag="CONVERSA"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, PANEL)
    tb(s, 0.75, 0.45, 4, 0.4, tag, size=13, bold=True, color=CYAN)
    tb(s, 0.85, 2.0, 8.3, 2.2, title, size=34, bold=True, color=WHITE)
    if hint:
        tb(s, 0.85, 4.8, 8.0, 1.5, hint, size=17, color=INK_SOFT, italic=True)
    bar(s)


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1
    title_slide(
        prs,
        "Fundamentos da BTECH",
        "Alinhamento com Bruno · Validar direção e seguir com o projeto.",
        dark=True,
    )

    # 2
    s = title_slide(prs, "Por que agora?", dark=False)
    bullets(
        s,
        0.85,
        2.7,
        8.3,
        4.0,
        [
            "O alinhamento de fundação importa — define para onde vamos.",
            "Mas não podemos perder o timing. É agora.",
            "",
            "Objetivo desta conversa:",
            "Validar o conceito em 45–60 min e voltar à execução.",
            "Não é workshop de grupo — é checagem de direção com quem aglutina o time.",
        ],
        size=20,
    )

    # 3 abertura
    question_slide(
        prs,
        "Por que a BTECH existe?",
        "Antes dos slides: qual é a sua leitura hoje?\n(5 min — você fala, eu escuto.)",
    )

    # 4 era
    s = title_slide(prs, "Uma mudança de era", dark=False)
    bullets(
        s,
        0.85,
        2.5,
        8.3,
        4.2,
        [
            "A Revolução Industrial não mudou só fábricas — reorganizou sociedade, trabalho, educação, cidades.",
            "",
            "IA, agentes e automação podem fazer algo parecido.",
            "",
            "Pergunta: estamos vendo nova tecnologia ou nova reorganização social?",
            "",
            "A IA deve ampliar possibilidades humanas — não substituí-las.",
        ],
        size=19,
    )

    # 5 futuros
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, LIGHT_BG)
    tb(s, 0.85, 0.5, 8.5, 0.7, "Dois futuros", size=30, bold=True, color=INK)
    for col, (tit, cor, items, x) in enumerate(
        [
            (
                "Concentração",
                RGBColor(35, 25, 55),
                ["Riqueza e decisão concentram", "Poucos controlam muito", "Autonomia diminui"],
                0.7,
            ),
            (
                "Ampliação",
                RGBColor(15, 40, 70),
                ["Acesso e oportunidade ampliam", "PMEs ganham capacidade", "Mais liberdade humana"],
                5.2,
            ),
        ]
    ):
        sh = s.shapes.add_shape(1, Inches(x), Inches(1.4), Inches(4.3), Inches(4.5))
        sh.fill.solid()
        sh.fill.fore_color.rgb = cor
        sh.line.fill.background()
        tb(s, x + 0.25, 1.6, 3.8, 0.5, tit, size=21, bold=True, color=WHITE)
        bullets(s, x + 0.25, 2.2, 3.8, 3.5, items, size=17, color=INK_SOFT)

    # 6 inversão
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, LIGHT_BG)
    tb(s, 0.85, 0.5, 8.5, 0.7, "A pergunta que nos diferencia", size=28, bold=True, color=INK)
    rows = [
        ("Mercado", "BTECH"),
        ("Como ganhamos dinheiro com IA?", "Como ampliamos capacidade humana e bem-estar com IA?"),
        ("Produtividade e custo", "Quem ganha, perde, fica de fora, passa a ter acesso"),
        ("Substituir pessoas", "Potencializar pessoas"),
    ]
    y = 1.5
    for a, b in rows:
        tb(s, 0.85, y, 4.0, 0.9, a, size=16, bold=(a in ("Mercado", "BTECH")), color=VIOLET if a == "BTECH" else INK)
        tb(s, 5.0, y, 4.5, 0.9, b, size=16, bold=(b == "BTECH"), color=VIOLET if b == "BTECH" else INK)
        y += 0.95
    tb(
        s,
        0.85,
        5.9,
        8.3,
        0.7,
        "Como garantir ampliação — e não concentração de poder?",
        size=17,
        bold=True,
        color=VIOLET,
    )

    # 7 prática
    s = title_slide(prs, "O que já provamos no chão", dark=False)
    tb(
        s,
        0.85,
        2.6,
        4.0,
        3.8,
        "Terapeuta + agente\n\nParece automação.\nÉ capacidade operacional de clínica maior — sem substituir o profissional.",
        size=18,
    )
    tb(
        s,
        5.1,
        2.6,
        4.0,
        3.8,
        "PMEs · Atmo · consultoria\n\nDistribuir acesso a cloud e IA.\nReduzir desigualdade tecnológica.",
        size=18,
    )
    tb(
        s,
        0.85,
        5.9,
        8.3,
        0.6,
        "Não é ONG — é infraestrutura de transição. Negócio sustentável.",
        size=16,
        bold=True,
        color=VIOLET,
    )

    # 8 B + proposta
    s = title_slide(prs, "O B e a proposta", dark=True)
    bullets(
        s,
        0.85,
        2.4,
        8.3,
        4.5,
        [
            "B = ampliação (quem opera, quem decide, quem acessa). Cloud e IA são ferramentas.",
            "",
            "PORQUÊ: prosperar na nova economia da IA — com mais autonomia e acesso.",
            "COMO: tecnologia complexa → solução acessível, perto da operação.",
            "O QUÊ hoje: cloud, backup, infra, agentes, automação.",
            "O QUÊ depois: educação, plataformas, produtos — o quê muda, o porquê fica.",
        ],
        size=18,
        color=INK_SOFT,
    )

    # 9 fechamento
    question_slide(
        prs,
        "Por que a BTECH existe?",
        "Fechamento · Valida, ajusta ou refuta?\nO que mudou depois desta conversa?\nSe alinhado: seguimos — timing é agora.",
        tag="VALIDAÇÃO",
    )

    # 10 próximo passo
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, INK)
    tb(s, 0.85, 2.2, 8.3, 1.0, "Próximo passo", size=32, bold=True, color=WHITE, center=True)
    tb(
        s,
        0.85,
        3.3,
        8.3,
        2.5,
        "Direção compartilhada → execução.\n\nRegistrar 1 frase de porquê (se fizer sentido).\nApresentar aos demais sócios quando for o momento.\n\nO futuro não é algo que encontramos.\nÉ algo que construímos.",
        size=22,
        color=INK_SOFT,
        center=True,
    )
    bar(s)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Gerado: {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
